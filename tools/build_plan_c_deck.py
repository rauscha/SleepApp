#!/usr/bin/env python3
"""Generate the Plan C (native rebuild) strategy deck.

Midnight Editorial Minimalism: deep ink background, warm stone text,
moonlit-sage accent, serif (Georgia) headings.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette ---------------------------------------------------------------
INK      = RGBColor(0x0E, 0x10, 0x16)  # deep background
PANEL    = RGBColor(0x17, 0x1B, 0x23)  # raised panel
PANEL2   = RGBColor(0x20, 0x25, 0x2E)
STONE    = RGBColor(0xE7, 0xE2, 0xD9)  # primary warm text
MUTE     = RGBColor(0x9A, 0x97, 0x8E)  # muted text
SAGE     = RGBColor(0x9D, 0xB8, 0xA6)  # moonlit sage accent
SAGE_DK  = RGBColor(0x5E, 0x72, 0x66)
EMBER    = RGBColor(0xCB, 0x86, 0x5E)  # warm warning accent
GREEN    = RGBColor(0x86, 0xB0, 0x8E)
HEAD = "Georgia"
BODY = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    return tb, tf


def setrun(r, text, size, color, font=BODY, bold=False, italic=False):
    r.text = text; f = r.font
    f.size = Pt(size); f.name = font; f.bold = bold; f.italic = italic
    f.color.rgb = color


def para(tf, text, size, color, font=BODY, bold=False, italic=False,
         space_after=6, space_before=0, level=0, bullet=None, align=None,
         first=False, line=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.space_after = Pt(space_after); p.space_before = Pt(space_before)
    p.level = level
    if align: p.alignment = align
    if line: p.line_spacing = line
    if bullet:
        r0 = p.add_run(); setrun(r0, bullet + "  ", size, SAGE, font, bold=True)
        r = p.add_run(); setrun(r, text, size, color, font, bold, italic)
    else:
        r = p.add_run(); setrun(r, text, size, color, font, bold, italic)
    return p


def accent_rule(s, x, y, w=Inches(1.1)):
    ln = s.shapes.add_shape(1, x, y, w, Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = SAGE; ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def header(s, kicker, title, title_size=34):
    box(s, Inches(0.7), Inches(0.42), Inches(11.9), Inches(0.3))[1]  # spacer
    _, tf = box(s, Inches(0.7), Inches(0.40), Inches(11.9), Inches(0.4))
    para(tf, kicker.upper(), 12, SAGE, BODY, bold=True, first=True)
    accent_rule(s, Inches(0.72), Inches(0.92))
    _, tf2 = box(s, Inches(0.7), Inches(1.02), Inches(11.9), Inches(1.0))
    para(tf2, title, title_size, STONE, HEAD, bold=True, first=True, line=1.04)


def panel(s, x, y, w, h, fill=PANEL):
    p = s.shapes.add_shape(1, x, y, w, h)
    p.fill.solid(); p.fill.fore_color.rgb = fill; p.line.color.rgb = PANEL2
    p.line.width = Pt(0.75); p.shadow.inherit = False
    return p


def pagenum(s, n, label):
    _, tf = box(s, Inches(0.7), Inches(7.04), Inches(11.9), Inches(0.34))
    p = para(tf, label, 9, MUTE, BODY, first=True)
    r = p.add_run(); setrun(r, "     ·     Sleep — Plan C", 9, SAGE_DK, BODY)
    _, tf2 = box(s, Inches(12.3), Inches(7.04), Inches(0.8), Inches(0.34))
    para(tf2, str(n), 9, MUTE, BODY, first=True, align=PP_ALIGN.RIGHT)


# ===========================================================================
# 1 — TITLE
s = slide()
accent_rule(s, Inches(0.9), Inches(2.55), Inches(1.4))
_, tf = box(s, Inches(0.85), Inches(2.7), Inches(11.6), Inches(2.4))
para(tf, "Going Native — Plan C", 52, STONE, HEAD, bold=True, first=True, line=1.0)
para(tf, "A ground-up native iOS + Android rebuild of Sleep,", 22, STONE, HEAD,
     italic=True, space_before=14, line=1.15)
para(tf, "driven by a long-running, self-checking multi-agent AI workflow.", 22,
     STONE, HEAD, italic=True, line=1.15)
_, tf2 = box(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.0))
para(tf2, "Feasibility  ·  pros & cons  ·  costs  ·  licensing (the big one)", 15,
     SAGE, BODY, first=True)
para(tf2, "Prepared overnight · 2026-06-15 · while the Howler html5 fix ships to your phone",
     12, MUTE, BODY, space_before=6)

# ===========================================================================
# 2 — TL;DR
s = slide(); header(s, "Bottom line up front", "The code is the easy part")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.6))
para(tf, "Engineering feasibility: HIGH.", 17, STONE, BODY, bold=True, first=True,
     bullet="▸", space_after=4)
para(tf, "A multi-agent Opus 4.8 / Fable 5 pipeline can realistically produce most of two native codebases — scene engine, UI, integrations, tests, CI.",
     14, MUTE, BODY, space_after=12)
para(tf, "The hard parts are NOT the code:", 17, STONE, BODY, bold=True, bullet="▸", space_after=4)
para(tf, "1)  Real-device overnight-audio QA — the exact thing emulators fake badly.   2)  App Store / Play review.   3)  Asset & dependency licensing — your biggest concern, and correctly so.",
     14, MUTE, BODY, space_after=12)
para(tf, "A bonus you'd get back:", 17, STONE, BODY, bold=True, bullet="▸", space_after=4)
para(tf, "Native sample-accurate looping RESTORES the Eno prime-offset engine and the synth-bed glue we just gave up in the web pivot.",
     14, MUTE, BODY, space_after=12)
para(tf, "My recommendation, up front:", 17, SAGE, BODY, bold=True, bullet="▸", space_after=4)
para(tf, "Test tonight's web fix first. If you still want native, pilot the cheap middle rung (native-audio shell) before committing to the full ground-up rewrite. Details on the last slide.",
     14, STONE, BODY)
pagenum(s, 2, "Executive summary")

# ===========================================================================
# 3 — WHY NATIVE
s = slide(); header(s, "The case for leaving the browser", "Why even consider native")
_, tf = box(s, Inches(0.7), Inches(2.15), Inches(12.0), Inches(4.6))
para(tf, "Tonight's Howler html5 fix is the best the browser allows — but the browser is always a guest in the OS's audio house.",
     15, STONE, BODY, first=True, bullet="—", space_after=12)
para(tf, "Native hands you the actual primitive the browser only imitates:", 15, STONE,
     BODY, bullet="—", space_after=6)
para(tf, "iOS:  AVAudioEngine / AVAudioSession (.playback) + the audio background mode — the OS owns the loop.",
     14, MUTE, BODY, level=1, space_after=4)
para(tf, "Android:  Media3 / ExoPlayer + a foreground media service — playback survives doze, screen-off, low memory.",
     14, MUTE, BODY, level=1, space_after=12)
para(tf, "This is precisely what Spotify, Calm, and YouTube ship. “Stay asleep all night” stops being a fight and becomes a guarantee.",
     15, SAGE, BODY, italic=True, bullet="—")
pagenum(s, 3, "Rationale")

# ===========================================================================
# 4 — SPECTRUM
s = slide(); header(s, "Don't skip the middle", "Three ways to “go native”")
cols = [
    ("RUNG 1", "Native-audio shell", "Capacitor / TWA wrapper around today's web UI + a native background-audio plugin.",
     "Effort: LOW  ·  reuse ~95%", "Robustness: HIGH (native audio)", "The 80/20. Pilot this first.", SAGE),
    ("RUNG 2", "Cross-platform rewrite", "One codebase in React Native / Flutter / Kotlin Multiplatform, with native audio modules.",
     "Effort: MEDIUM–HIGH", "Robustness: HIGH", "One team, two stores.", STONE),
    ("RUNG 3", "Ground-up native", "Separate Swift (iOS) + Kotlin (Android), hand-built per platform. WHAT YOU ASKED ABOUT.",
     "Effort: HIGHEST", "Robustness: HIGHEST", "Best fidelity & longevity.", EMBER),
]
x = Inches(0.7); w = Inches(3.86); gap = Inches(0.18); y = Inches(2.25); h = Inches(4.1)
for i, (k, t, d, e, r2, note, c) in enumerate(cols):
    px = Emu(int(x) + i * (int(w) + int(gap)))
    panel(s, px, y, w, h)
    _, tf = box(s, Emu(int(px) + Emu(0.22*914400)), Emu(int(y)+Emu(0.2*914400)),
                Emu(int(w) - Emu(0.44*914400)), Emu(int(h)-Emu(0.4*914400)))
    para(tf, k, 12, c, BODY, bold=True, first=True, space_after=2)
    para(tf, t, 19, STONE, HEAD, bold=True, space_after=8, line=1.05)
    para(tf, d, 12.5, MUTE, BODY, space_after=10, line=1.16)
    para(tf, e, 12, STONE, BODY, bold=True, space_after=2)
    para(tf, r2, 12, STONE, BODY, bold=True, space_after=8)
    para(tf, note, 12, c, BODY, italic=True)
_, tf = box(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.5))
para(tf, "You asked specifically about Rung 3 — the rest of the deck focuses there. But I'd pilot Rung 1 first: it validates OS background audio for days of effort, not months.",
     12.5, SAGE, BODY, italic=True, first=True)
pagenum(s, 4, "Option spectrum")

# ===========================================================================
# 5 — ARCHITECTURE
s = slide(); header(s, "Plan C, concretely", "The ground-up native architecture")
panel(s, Inches(0.7), Inches(2.15), Inches(5.85), Inches(3.5))
_, tf = box(s, Inches(0.95), Inches(2.35), Inches(5.4), Inches(3.2))
para(tf, "iOS — Swift", 18, SAGE, HEAD, bold=True, first=True, space_after=6)
para(tf, "AVAudioEngine graph; one AVAudioPlayerNode per layer scheduling looping buffers at sample-accurate prime offsets.",
     13, STONE, BODY, bullet="·", space_after=5, line=1.15)
para(tf, "AVAudioSession .playback + UIBackgroundModes audio.", 13, STONE, BODY, bullet="·", space_after=5)
para(tf, "MPNowPlayingInfoCenter + remote command center for lock-screen transport.",
     13, STONE, BODY, bullet="·", space_after=5)
para(tf, "MPVolumeView / fade ramps for the sleep timer.", 13, STONE, BODY, bullet="·")
panel(s, Inches(6.75), Inches(2.15), Inches(5.85), Inches(3.5))
_, tf = box(s, Inches(7.0), Inches(2.35), Inches(5.4), Inches(3.2))
para(tf, "Android — Kotlin", 18, SAGE, HEAD, bold=True, first=True, space_after=6)
para(tf, "Media3 / ExoPlayer looping sources per layer (or AudioTrack for full control).",
     13, STONE, BODY, bullet="·", space_after=5, line=1.15)
para(tf, "MediaSessionService + foreground service, type = mediaPlayback.", 13, STONE, BODY, bullet="·", space_after=5)
para(tf, "Survives Doze, screen-off, low-memory reclaim.", 13, STONE, BODY, bullet="·", space_after=5)
para(tf, "AudioFocus handling for headset / call interruptions.", 13, STONE, BODY, bullet="·")
_, tf = box(s, Inches(0.7), Inches(5.85), Inches(12.0), Inches(1.2))
para(tf, "Shared & carried over unchanged:  the scene JSON format, the design system, story text (Claude) + narration (ElevenLabs).",
     13.5, STONE, BODY, first=True, bullet="▸", space_after=6)
para(tf, "The headline win: native sample-accurate looping brings the prime-offset Eno engine AND the synth bed back — the two things the web pivot had to sacrifice.",
     13.5, SAGE, BODY, italic=True, bullet="▸")
pagenum(s, 5, "Architecture")

# ===========================================================================
# 6 — MULTIAGENT WORKFLOW
s = slide(); header(s, "How the build actually runs", "The long-running multi-agent workflow")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(6.1), Inches(4.7))
para(tf, "Orchestrator agent", 16, SAGE, BODY, bold=True, first=True, bullet="◆", space_after=3)
para(tf, "owns the plan + backlog, checkpoints to git, opens PRs at milestones, and spawns specialists:",
     13, MUTE, BODY, space_after=8, line=1.15)
for a, d in [
    ("Architect", "system design, scene-engine port"),
    ("iOS-audio / Android-audio", "the overnight-survival core"),
    ("UI / design-system", "screens from the brand spec"),
    ("Integrations", "Claude text + ElevenLabs TTS"),
    ("Test / QA", "unit, snapshot, integration"),
    ("Reviewer / critic", "diff critique + license scan"),
]:
    p = para(tf, a + " — ", 13, STONE, BODY, bold=True, bullet="·", space_after=3)
    r = p.add_run(); setrun(r, d, 13, MUTE, BODY)
panel(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(4.55))
_, tf = box(s, Inches(7.25), Inches(2.3), Inches(5.1), Inches(4.2))
para(tf, "Self-checking loops", 16, SAGE, BODY, bold=True, first=True, space_after=6)
para(tf, "Every change gated on: compile → unit tests → UI snapshot → lint. The Reviewer agent critiques each diff; nightly integration runs on CI (macOS runners for iOS, Android emulator).",
     13, STONE, BODY, space_after=10, line=1.18)
para(tf, "Runs unattended for days/weeks", 16, SAGE, BODY, bold=True, space_after=6)
para(tf, "You review at milestone PRs, not line by line.", 13, STONE, BODY, space_after=10, line=1.18)
para(tf, "Model split", 16, SAGE, BODY, bold=True, space_after=6)
para(tf, "Opus 4.8 plans, owns the audio core, and gates merges. Fable 5 does the high-throughput implementation passes. (Appendix.)",
     13, STONE, BODY, line=1.18)
pagenum(s, 6, "Workflow")

# ===========================================================================
# 7 — AGENTS CAN / CAN'T
s = slide(); header(s, "An honest division of labour", "What the agents do vs. what still needs you")
panel(s, Inches(0.7), Inches(2.15), Inches(5.85), Inches(4.4))
_, tf = box(s, Inches(0.95), Inches(2.35), Inches(5.4), Inches(4.0))
para(tf, "Agents do this well", 18, GREEN, HEAD, bold=True, first=True, space_after=8)
for t in ["Two native codebases — engine, UI, glue",
          "Porting the scene JSON format verbatim",
          "Test suites + CI pipelines",
          "UI built from the design system spec",
          "Refactors, docs, dependency wiring",
          "License-scanning dependencies"]:
    para(tf, t, 13.5, STONE, BODY, bullet="✓", space_after=6, line=1.12)
panel(s, Inches(6.75), Inches(2.15), Inches(5.85), Inches(4.4))
_, tf = box(s, Inches(7.0), Inches(2.35), Inches(5.4), Inches(4.0))
para(tf, "Still needs a human", 18, EMBER, HEAD, bold=True, first=True, space_after=8)
for t in ["Real-device overnight audio QA — emulators don't reproduce OS throttling faithfully",
          "App Store / Play submission, signing, certs",
          "Responding to review rejections",
          "Licensing decisions & asset clearance (judgment + contracts)",
          "Developer accounts, banking, tax setup"]:
    para(tf, t, 13.5, STONE, BODY, bullet="•", space_after=6, line=1.12)
pagenum(s, 7, "Feasibility")

# ===========================================================================
# 8 — PROS
s = slide(); header(s, "The upside", "Pros of the native rebuild")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.8))
for t, d in [
    ("Bulletproof overnight audio", "The OS owns the loop. The thing we've fought for a week stops being a risk."),
    ("The Eno engine comes back", "Sample-accurate native looping restores the prime offsets and the synth-bed glue."),
    ("App Store / Play presence", "Discoverability, credibility, real install funnel, optional paid/IAP later."),
    ("Better performance & offline", "No browser tax; assets bundled; instant cold start; true offline."),
    ("Clean privacy story", "No accounts, no telemetry → trivially clean App Privacy labels (a review advantage)."),
    ("Platform features", "Lock-screen art, Siri/Assistant, widgets, Focus/Bedtime integration, CarPlay/AA later."),
]:
    p = para(tf, t + "  —  ", 15, SAGE, BODY, bold=True, first=(t.startswith("Bullet")),
             bullet="▸", space_after=9)
    r = p.add_run(); setrun(r, d, 14, STONE, BODY)
pagenum(s, 8, "Pros")

# ===========================================================================
# 9 — CONS
s = slide(); header(s, "The downside", "Cons & risks")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.8))
for t, d in [
    ("You lose instant web deploy", "No more “merge → live on your phone in 2 minutes.” Builds, signing, store review now gate every release."),
    ("Two platforms to maintain", "(Rung 3) Swift + Kotlin = double the surface, double the OS-update churn, indefinitely."),
    ("App review latency & rejection", "Days per cycle; background-audio + any AI content draw scrutiny."),
    ("Device-QA can't be fully automated", "The one failure mode that matters most still needs real phones in real overnight conditions."),
    ("Multi-agent drift & cost", "Long autonomous runs can wander; needs the critic loop + your milestone reviews to stay on-brief."),
    ("Up-front cost & time", "Weeks-to-months and real dollars before the first install — vs. the web app that already works."),
]:
    p = para(tf, t + "  —  ", 15, EMBER, BODY, bold=True, first=(t.startswith("You lose")),
             bullet="▸", space_after=9)
    r = p.add_run(); setrun(r, d, 14, STONE, BODY)
pagenum(s, 9, "Cons")

# ===========================================================================
# 10 — COSTS
s = slide(); header(s, "What it costs", "Budget — order-of-magnitude")
rows = [
    ("AI build (API tokens)", "One-time", "Low hundreds → low thousands of $",
     "A 2-platform ground-up app is tens of millions of tokens over weeks. Fable-5-heavy bulk + Opus for core/review keeps it toward the low end; all-Opus pushes it up."),
    ("Your review time", "One-time", "The real cost", "Milestone reviews, device testing, decisions. Hours, not dollars — but the binding constraint."),
    ("Apple Developer", "$99 / year", "Required", "Per Apple account, to ship on the App Store."),
    ("Google Play", "$25 once", "Required", "One-time registration."),
    ("ElevenLabs", "Subscription", "Paid tier required", "Commercial ownership of narration needs a paid plan (free tier = no commercial use)."),
    ("Anthropic API", "Usage", "Per story", "Story-text generation (already in use)."),
    ("Server (optional)", "Monthly", "Only if moved", "If TTS/LLM move server-side off the client; otherwise $0."),
    ("Maintenance", "Ongoing", "Annual", "OS updates, cert renewals, store-policy changes × two codebases."),
]
y = Inches(2.2); rh = Inches(0.55); x = Inches(0.7)
cw = [Inches(2.7), Inches(1.5), Inches(2.2), Inches(5.5)]
hdr = ["Item", "Cadence", "Estimate", "Notes"]
# header row
panel(s, x, y, Inches(11.9), Inches(0.42), PANEL2)
cx = x
for i, htxt in enumerate(hdr):
    _, tf = box(s, Emu(int(cx)+Emu(0.1*914400)), Emu(int(y)+Emu(0.04*914400)), cw[i], Inches(0.34))
    para(tf, htxt, 12.5, SAGE, BODY, bold=True, first=True)
    cx = Emu(int(cx) + int(cw[i]))
yy = Emu(int(y) + int(Inches(0.42)))
for ri, row in enumerate(rows):
    if ri % 2 == 0:
        panel(s, x, yy, Inches(11.9), rh, PANEL)
    cx = x
    for i, cell in enumerate(row):
        _, tf = box(s, Emu(int(cx)+Emu(0.1*914400)), Emu(int(yy)+Emu(0.03*914400)), cw[i], rh)
        col = STONE if i == 0 else (SAGE if i == 2 else MUTE)
        bold = (i == 0)
        para(tf, cell, 10.5 if i == 3 else 11.5, col, BODY, bold=bold, first=True, line=1.04)
        cx = Emu(int(cx) + int(cw[i]))
    yy = Emu(int(yy) + int(rh))
pagenum(s, 10, "Costs")

# ===========================================================================
# 11 — TIMELINE
s = slide(); header(s, "How long", "Period — phased, estimate")
phases = [
    ("Pilot — Rung 1 audio shell", "~ days", "Wrap today's UI, drop in a native background-audio plugin. Proves OS overnight audio cheaply, before any rewrite spend."),
    ("Plan C MVP — one platform", "~ 2–4 weeks", "Core scenes + sleep timer + Night Drift on iOS (or Android). Stories deferred. Agent build + your milestone reviews."),
    ("Second platform + parity", "~ +2–3 weeks", "Port to the other OS; shared scene format makes this faster than the first."),
    ("Stories, polish, submission", "~ +3–6 weeks", "TTS/LLM integration, design polish, store assets, privacy labels, and the review cycles themselves."),
]
y = Inches(2.25)
for t, dur, d in phases:
    panel(s, Inches(0.7), y, Inches(11.9), Inches(1.02))
    _, tf = box(s, Inches(0.95), Emu(int(y)+Emu(0.14*914400)), Inches(8.3), Inches(0.8))
    para(tf, t, 16, STONE, HEAD, bold=True, first=True, space_after=2)
    para(tf, d, 12.5, MUTE, BODY, line=1.12)
    _, tf2 = box(s, Inches(9.4), Emu(int(y)+Emu(0.28*914400)), Inches(3.0), Inches(0.5))
    para(tf2, dur, 18, SAGE, HEAD, bold=True, first=True, align=PP_ALIGN.RIGHT)
    y = Emu(int(y) + int(Inches(1.12)))
_, tf = box(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.4))
para(tf, "Estimates dominated by device-QA and store-review latency — not code-writing speed.",
     12, SAGE, BODY, italic=True, first=True)
pagenum(s, 11, "Timeline")

# ===========================================================================
# 12 — LICENSING OVERVIEW
s = slide(); header(s, "Your biggest concern", "Licensing — the seven buckets")
buckets = [
    ("Audio — Pixabay scenes", GREEN, "Low risk", "Commercial OK, no attribution; bundling in an app is fine."),
    ("Audio — freetousesounds", EMBER, "Action needed", "Requires attribution + bars redistributing originals. forest-evening uses it."),
    ("AI narration — ElevenLabs", GREEN, "OK if paid tier", "Paid plan = you own output, perpetual commercial rights."),
    ("AI story text — Claude", GREEN, "OK", "Outputs assigned to you under commercial terms."),
    ("Fonts", EMBER, "Verify", "Editorial serif needs an APP-EMBEDDING license; default to OFL."),
    ("Photography", EMBER, "Open gap", "Scene photos need commercial, app-distributable image licenses."),
    ("Code & OSS deps", GREEN, "OK", "AI output rights are yours; deps are permissive (MIT/Apache)."),
]
y = Inches(2.2); rh = Inches(0.62)
for i, (t, c, status, d) in enumerate(buckets):
    panel(s, Inches(0.7), y, Inches(11.9), Inches(0.54))
    chip = s.shapes.add_shape(1, Inches(0.7), y, Inches(0.12), Inches(0.54))
    chip.fill.solid(); chip.fill.fore_color.rgb = c; chip.line.fill.background(); chip.shadow.inherit = False
    _, tf = box(s, Inches(0.95), Emu(int(y)+Emu(0.06*914400)), Inches(3.7), Inches(0.45))
    para(tf, t, 13.5, STONE, BODY, bold=True, first=True)
    _, tf2 = box(s, Inches(4.7), Emu(int(y)+Emu(0.06*914400)), Inches(1.9), Inches(0.45))
    para(tf2, status, 12.5, c, BODY, bold=True, first=True)
    _, tf3 = box(s, Inches(6.6), Emu(int(y)+Emu(0.07*914400)), Inches(6.0), Inches(0.45))
    para(tf3, d, 11.5, MUTE, BODY, first=True, line=1.04)
    y = Emu(int(y) + int(rh))
pagenum(s, 12, "Licensing — overview")

# ===========================================================================
# 13 — LICENSING: AUDIO
s = slide(); header(s, "Licensing in detail · 1 of 3", "Audio assets — the live one")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.8))
para(tf, "Pixabay (most scenes — rain, monsoon, ocean, etc.)", 15, GREEN, BODY, bold=True,
     first=True, bullet="✓", space_after=3)
para(tf, "Pixabay Content License: free, commercial use OK, no attribution. Bundling inside the app is fine. You only can't sell the sound files as a standalone pack. Lowest-risk bucket.",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "freetousesounds (forest-evening: forest-rain + creek variants)", 15, EMBER, BODY,
     bold=True, bullet="!", space_after=3)
para(tf, "Their license requires ATTRIBUTION (a visible credit/link) and prohibits redistributing the original sounds unmodified — and flags that app-developer commercial use may need a separate non-exclusive license. Our files are re-encoded/trimmed, which helps, but “modified” should mean genuinely redesigned.",
     13, STONE, BODY, space_after=4, line=1.16)
para(tf, "Reconcile with “no chrome”: a tucked-away About/Credits screen is the usual fix — a media-credits list isn't marketing UI. Or replace these specific clips with Pixabay/OFL-clean equivalents.",
     13, SAGE, BODY, italic=True, level=1, space_after=12)
para(tf, "“User-provided” & any unknown-source files (e.g. wind-2)", 15, EMBER, BODY, bold=True,
     bullet="!", space_after=3)
para(tf, "Provenance must be confirmed before commercial distribution. The Reviewer agent can produce a per-file license manifest from the existing .json sidecars as step one.",
     13, STONE, BODY, line=1.16)
pagenum(s, 13, "Licensing — audio")

# ===========================================================================
# 14 — LICENSING: AI CONTENT
s = slide(); header(s, "Licensing in detail · 2 of 3", "AI-generated content")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.8))
para(tf, "ElevenLabs — narration (TTS)", 15, GREEN, BODY, bold=True, first=True, bullet="✓", space_after=3)
para(tf, "On a PAID plan you own the generated audio with perpetual commercial rights — even after cancelling. Free tier = no commercial use + forced “elevenlabs.io” attribution. They retain a license to train on your content. Their Sound-FX / Music products have standalone-redistribution limits, but in-app narration is squarely fine.",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "Anthropic / Claude — story text", 15, GREEN, BODY, bold=True, bullet="✓", space_after=3)
para(tf, "Under the commercial terms, rights to outputs are assigned to you; usable in-app. Keep generation policy-compliant (the Prohibited Use Policy).",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "If you ever add AI-generated MUSIC beds", 15, EMBER, BODY, bold=True, bullet="!", space_after=3)
para(tf, "Check that specific product's commercial clearance and any marketplace redistribution limits before shipping it as a scene layer. Music rights are materially trickier than TTS.",
     13, STONE, BODY, line=1.16)
para(tf, "Net: the AI-content buckets are the SAFE ones — provided narration runs on a paid ElevenLabs tier.",
     13.5, SAGE, BODY, italic=True, space_before=10, bullet="▸")
pagenum(s, 14, "Licensing — AI content")

# ===========================================================================
# 15 — LICENSING: FONTS/PHOTOS/CODE
s = slide(); header(s, "Licensing in detail · 3 of 3", "Fonts, photos, code & deps")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(4.8))
para(tf, "Fonts — the editorial serif", 15, EMBER, BODY, bold=True, first=True, bullet="!", space_after=3)
para(tf, "Many foundries license web use ≠ app embedding. Confirm an app-embedding/redistribution license, or default to an SIL-OFL serif (e.g. a Google Fonts serif) which is safe to embed and ship. Lowest-friction path.",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "Photography — biggest open visual gap", 15, EMBER, BODY, bold=True, bullet="!", space_after=3)
para(tf, "Already flagged in the brief. Store distribution needs commercial, app-distributable image licenses (model/property releases where relevant) — or shoot/replace. Don't ship borrowed imagery.",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "Code written by Opus / Fable", 15, GREEN, BODY, bold=True, bullet="✓", space_after=3)
para(tf, "Anthropic grants you rights to model outputs — the generated code is yours. No copyleft contamination as long as agents don't paste GPL/AGPL snippets; the Reviewer agent runs a license scan to enforce that.",
     13, STONE, BODY, space_after=12, line=1.16)
para(tf, "Open-source dependencies", 15, GREEN, BODY, bold=True, bullet="✓", space_after=3)
para(tf, "ExoPlayer (Apache-2), Howler (MIT), RN/Flutter (permissive). Ship a NOTICES file; no obligations beyond attribution-in-app.",
     13, STONE, BODY, line=1.16)
pagenum(s, 15, "Licensing — fonts/photos/code")

# ===========================================================================
# 16 — ENTITLEMENTS & TENSION
s = slide(); header(s, "One thing to decide early", "Background entitlements & the “no notifications” tension")
panel(s, Inches(0.7), Inches(2.15), Inches(5.85), Inches(3.0))
_, tf = box(s, Inches(0.95), Inches(2.35), Inches(5.4), Inches(2.7))
para(tf, "iOS", 17, SAGE, HEAD, bold=True, first=True, space_after=6)
para(tf, "The “audio” background mode is free but reviewed — a sleep app legitimately qualifies. The session must keep producing audible output or the OS ends it (no silent keep-alive tricks needed; real audio is playing).",
     13, STONE, BODY, line=1.18)
panel(s, Inches(6.75), Inches(2.15), Inches(5.85), Inches(3.0))
_, tf = box(s, Inches(7.0), Inches(2.35), Inches(5.4), Inches(2.7))
para(tf, "Android 13+", 17, EMBER, HEAD, bold=True, first=True, space_after=6)
para(tf, "Reliable background playback REQUIRES a foreground media service — which mandates an ongoing media notification. This rubs against the brief's “no notifications, ever.”",
     13, STONE, BODY, line=1.18)
_, tf = box(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.6))
para(tf, "Your call:  a media-playback notification is an OS transport control (play/pause on the lock screen), not a push or marketing alert. It's arguably in-spirit — and it's the price of bulletproof Android background audio. Worth deciding before the rebuild, because it shapes the Android architecture.",
     13.5, SAGE, BODY, italic=True, first=True, bullet="▸", line=1.2)
para(tf, "Bonus: the no-accounts / no-telemetry stance makes App Privacy labels trivially clean — a genuine advantage at review.",
     13, STONE, BODY, bullet="▸", space_before=6)
pagenum(s, 16, "Entitlements")

# ===========================================================================
# 17 — RISK MATRIX
s = slide(); header(s, "At a glance", "Risk matrix")
risks = [
    ("Device audio QA can't be fully automated", "High", "High", "Real-phone overnight test gates each milestone; consider a small device set."),
    ("App-store review (audio + AI content)", "Med", "Med", "Clean privacy labels; justify the audio entitlement; budget review cycles."),
    ("Audio licensing (freetousesounds / unknown)", "Med", "High", "Per-file manifest; add credits screen or replace clips."),
    ("Font / photo licensing", "Med", "Med", "Default to OFL fonts; source or shoot licensed imagery."),
    ("Multi-agent drift over a long run", "Med", "Med", "Critic loop + milestone PR reviews + tight scene-format contract tests."),
    ("Cost / time overrun", "Med", "Med", "Pilot Rung 1 first; Fable-heavy bulk; hard budget ceiling."),
    ("Two-codebase maintenance burden", "High", "Med", "Consider Rung 2 (one cross-platform codebase) instead of Rung 3."),
]
y = Inches(2.15); rh = Inches(0.6)
panel(s, Inches(0.7), y, Inches(11.9), Inches(0.4), PANEL2)
for cx, w2, lbl in [(Inches(0.8), Inches(5.6), "Risk"), (Inches(6.5), Inches(1.2), "Likely"),
                    (Inches(7.7), Inches(1.2), "Impact"), (Inches(8.95), Inches(3.6), "Mitigation")]:
    _, tf = box(s, cx, Emu(int(y)+Emu(0.03*914400)), w2, Inches(0.34))
    para(tf, lbl, 12, SAGE, BODY, bold=True, first=True)
yy = Emu(int(y)+int(Inches(0.4)))
def sev_color(v): return EMBER if v == "High" else (SAGE if v == "Med" else GREEN)
for i, (t, lk, im, mit) in enumerate(risks):
    if i % 2 == 0: panel(s, Inches(0.7), yy, Inches(11.9), rh, PANEL)
    _, tf = box(s, Inches(0.8), Emu(int(yy)+Emu(0.05*914400)), Inches(5.6), rh)
    para(tf, t, 11.5, STONE, BODY, bold=True, first=True, line=1.05)
    _, tf = box(s, Inches(6.5), Emu(int(yy)+Emu(0.05*914400)), Inches(1.2), rh)
    para(tf, lk, 11.5, sev_color(lk), BODY, bold=True, first=True)
    _, tf = box(s, Inches(7.7), Emu(int(yy)+Emu(0.05*914400)), Inches(1.2), rh)
    para(tf, im, 11.5, sev_color(im), BODY, bold=True, first=True)
    _, tf = box(s, Inches(8.95), Emu(int(yy)+Emu(0.06*914400)), Inches(3.6), rh)
    para(tf, mit, 10.5, MUTE, BODY, first=True, line=1.04)
    yy = Emu(int(yy)+int(rh))
pagenum(s, 17, "Risk")

# ===========================================================================
# 18 — RECOMMENDATION
s = slide(); header(s, "What I'd do", "Recommendation & decisions I need from you")
_, tf = box(s, Inches(0.7), Inches(2.1), Inches(12.0), Inches(2.5))
para(tf, "1.  Test tonight's web fix first.", 15, SAGE, BODY, bold=True, first=True, bullet="", space_after=3)
para(tf, "If the Howler html5 build survives the night, native becomes a reach-for-the-stars upgrade, not a rescue — which changes the math.",
     13, STONE, BODY, space_after=8, line=1.14)
para(tf, "2.  If you still want native, pilot Rung 1 (native-audio shell).", 15, SAGE, BODY, bold=True, bullet="", space_after=3)
para(tf, "Days of effort to validate real OS background audio on your phone before committing months to a ground-up rewrite.",
     13, STONE, BODY, space_after=8, line=1.14)
para(tf, "3.  Greenlight full Plan C only once the pilot proves the audio and the licensing checklist is clear.",
     15, SAGE, BODY, bold=True, bullet="", line=1.14)
panel(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(2.0))
_, tf = box(s, Inches(0.95), Inches(5.0), Inches(11.4), Inches(1.8))
para(tf, "Decisions I need from you", 15, EMBER, BODY, bold=True, first=True, space_after=6)
for t in ["Cross-platform (Rung 2, one codebase) vs. separate native (Rung 3)?",
          "The Android media-notification call — acceptable, or a dealbreaker?",
          "Budget ceiling for the agent build run?",
          "Story TTS/LLM — stay client-side, or move server-side?",
          "Asset clearance go/no-go: replace freetousesounds + source licensed photos?"]:
    para(tf, t, 12.5, STONE, BODY, bullet="•", space_after=3, line=1.08)
pagenum(s, 18, "Recommendation")

# ===========================================================================
# 19 — APPENDIX MODEL CHOICE
s = slide(); header(s, "Appendix", "Opus 4.8 vs. Fable 5 for the build")
panel(s, Inches(0.7), Inches(2.2), Inches(5.85), Inches(3.8))
_, tf = box(s, Inches(0.95), Inches(2.4), Inches(5.4), Inches(3.5))
para(tf, "Opus 4.8", 20, SAGE, HEAD, bold=True, first=True, space_after=8)
para(tf, "Deepest reasoning. Use for:", 13.5, STONE, BODY, space_after=6)
for t in ["System architecture", "The audio-engine core (the hard part)", "Security & licensing review", "The Reviewer/critic agent", "Merge gating"]:
    para(tf, t, 13, STONE, BODY, bullet="·", space_after=4)
panel(s, Inches(6.75), Inches(2.2), Inches(5.85), Inches(3.8))
_, tf = box(s, Inches(7.0), Inches(2.4), Inches(5.4), Inches(3.5))
para(tf, "Fable 5", 20, SAGE, HEAD, bold=True, first=True, space_after=8)
para(tf, "Fast, high throughput. Use for:", 13.5, STONE, BODY, space_after=6)
for t in ["Bulk implementation passes", "Test generation", "Boilerplate & glue code", "Docs & changelogs", "Routine refactors"]:
    para(tf, t, 13, STONE, BODY, bullet="·", space_after=4)
_, tf = box(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7))
para(tf, "Pattern: Opus plans and reviews; Fable executes; Opus gates merges. Cheaper than all-Opus, higher quality than all-Fable.",
     13.5, SAGE, BODY, italic=True, first=True, align=PP_ALIGN.CENTER)
pagenum(s, 19, "Appendix")

out = "/home/user/SleepApp/Plan-C-Native-Rebuild.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
